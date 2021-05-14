#pip3 install python-pptx -i https://pypi.doubanio.com/simple/ --trusted-host pypi.doubanio.com
# 加载库
import  os
from pptx import Presentation
from pptx.util import Cm, Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement


# 设置路径
work_path = './'
os.chdir(work_path)

# 实例化 ppt 文档对象
prs = Presentation()

# 插入幻灯片
blank_slide = prs.slide_layouts[6]
slide_1 = prs.slides.add_slide(blank_slide)

# 添加图片
img_name1  = 'test.jpg'
img_name  = 'test满月.png'
## 中央主图
pic = slide_1.shapes.add_picture(image_file=img_name1,
                           left=Inches(0),
                           top=Inches(0),
                           width=Inches(10),
                           height=Inches(5)
                          )
## 左上角小图
slide_1.shapes.add_picture(image_file=img_name,
                           left=Inches(0.5),
                           top=Inches(0.5),
                           width=Inches(3),
                           height=Inches(1.5)
                          )
## 右上角小图
slide_1.shapes.add_picture(image_file=img_name,
                           left=Inches(7.5),
                           top=Inches(0.5),
                           width=Inches(3),
                           height=Inches(1.5)
                          )
## 左下角小图
slide_1.shapes.add_picture(image_file=img_name,
                           left=Inches(0.5),
                           top=Inches(5),
                           width=Inches(3),
                           height=Inches(1.5)
                          )
# 右下角小图
slide_1.shapes.add_picture(image_file=img_name,
                           left=Inches(7.5),
                           top=Inches(5),
                           width=Inches(3),
                           height=Inches(1.5)
                          )
# 添加文本框
textbox= slide_1.shapes.add_textbox(left=Inches(1),
                                    top=Inches(6),
                                    width=Inches(8),
                                    height=Inches(2)
                                   )
## 向文本框加入文字
tf = textbox.text_frame
para = tf.add_paragraph()    # 添加段落
para.text = '神马都是浮云！！！'
para.alignment = PP_ALIGN.CENTER  # 居中
## 设置字体
font = para.font
font.size = Pt(36)    # 大小
font.name = '微软雅黑'    # 字体
font.bold = True    # 加粗
font.italic = True  # 倾斜
font.alpha = 0.30

font.color.rgb = RGBColor(225, 225, 0)  # 黄色

# 保存 ppt
prs.save('temp/test.pptx')